"""
党课汇报 PPT 生成脚本
《坚守正确政绩观，筑牢民航初心使命》
基于民航运输领域典型案例的党课讲稿改编而成
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree

# ============ 颜色配置（党政红主题）============
COLOR_PRIMARY = RGBColor(0xA8, 0x1F, 0x1F)      # 党政深红
COLOR_PRIMARY_DARK = RGBColor(0x7A, 0x14, 0x14)  # 暗红
COLOR_ACCENT = RGBColor(0xD4, 0xAF, 0x37)        # 金色
COLOR_DARK = RGBColor(0x1F, 0x2B, 0x3D)          # 深蓝灰
COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)          # 正文深灰
COLOR_TEXT_LIGHT = RGBColor(0x66, 0x66, 0x66)    # 浅灰说明文字
COLOR_BG_LIGHT = RGBColor(0xF5, 0xF1, 0xEC)      # 米白背景
COLOR_BG_CARD = RGBColor(0xFA, 0xF6, 0xF0)       # 卡片背景
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_GREEN = RGBColor(0x2E, 0x7D, 0x32)         # 正面示例绿
COLOR_RED_NEG = RGBColor(0xC6, 0x28, 0x28)       # 反面警示红

FONT_TITLE = "微软雅黑"
FONT_BODY = "微软雅黑"
FONT_CN_HEAVY = "微软雅黑"

# 16:9 宽屏
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_run(run, text, size=18, bold=False, color=None, font=FONT_BODY):
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    # 中文字体设置
    rPr = run._r.get_or_add_rPr()
    rFonts = rPr.find(qn('a:rFonts'))
    if rFonts is None:
        rFonts = etree.SubElement(rPr, qn('a:rFonts'))
    rFonts.set('eastAsia', font)
    rFonts.set('ascii', font)
    rFonts.set('hAnsi', font)


def add_text_box(slide, left, top, width, height, text, size=18, bold=False,
                 color=COLOR_TEXT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                 font=FONT_BODY, line_spacing=1.3):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    set_run(run, text, size=size, bold=bold, color=color, font=font)
    return tb


def add_multi_paragraphs(slide, left, top, width, height, paragraphs,
                         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                         line_spacing=1.4):
    """paragraphs: list of dicts with keys text, size, bold, color, font, align"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    for i, para in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = para.get('align', align)
        p.line_spacing = para.get('line_spacing', line_spacing)
        if 'space_before' in para:
            p.space_before = Pt(para['space_before'])
        if 'space_after' in para:
            p.space_after = Pt(para['space_after'])
        run = p.add_run()
        set_run(run,
                para['text'],
                size=para.get('size', 16),
                bold=para.get('bold', False),
                color=para.get('color', COLOR_TEXT),
                font=para.get('font', FONT_BODY))
    return tb


def add_rect(slide, left, top, width, height, fill_color,
             line_color=None, shadow=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
    if not shadow:
        shp.shadow.inherit = False
    shp.text_frame.text = ""
    return shp


def add_round_rect(slide, left, top, width, height, fill_color,
                   line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 left, top, width, height)
    shp.adjustments[0] = 0.12
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
    shp.shadow.inherit = False
    shp.text_frame.text = ""
    return shp


def add_page_footer(slide, page_num, total, section=""):
    # 底部装饰条
    add_rect(slide, Emu(0), SLIDE_H - Inches(0.35),
             SLIDE_W, Inches(0.35), COLOR_PRIMARY)
    # 左下：章节
    if section:
        add_text_box(slide, Inches(0.4), SLIDE_H - Inches(0.32),
                     Inches(8), Inches(0.3), section,
                     size=10, color=COLOR_WHITE, font=FONT_BODY,
                     anchor=MSO_ANCHOR.MIDDLE)
    # 右下：页码
    add_text_box(slide, SLIDE_W - Inches(1.5), SLIDE_H - Inches(0.32),
                 Inches(1.1), Inches(0.3),
                 f"{page_num} / {total}",
                 size=10, color=COLOR_WHITE, font=FONT_BODY,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def add_slide_header(slide, title, subtitle=""):
    # 左侧红色装饰条
    add_rect(slide, Emu(0), Inches(0.55), Inches(0.18), Inches(0.7),
             COLOR_PRIMARY)
    # 标题
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(11), Inches(0.55),
                 title, size=26, bold=True, color=COLOR_DARK,
                 anchor=MSO_ANCHOR.MIDDLE, font=FONT_TITLE)
    if subtitle:
        add_text_box(slide, Inches(0.5), Inches(1.05), Inches(11),
                     Inches(0.35), subtitle,
                     size=12, color=COLOR_TEXT_LIGHT,
                     font=FONT_BODY)
    # 标题下分隔线
    add_rect(slide, Inches(0.5), Inches(1.45),
             Inches(12.3), Emu(15000), COLOR_PRIMARY)
    add_rect(slide, Inches(0.5), Inches(1.47),
             Inches(0.8), Emu(30000), COLOR_ACCENT)


# ============ 开始构建幻灯片 ============
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]

TOTAL_PAGES = 17  # 总页数


# ---------------- Slide 1: 封面 ----------------
def build_cover():
    slide = prs.slides.add_slide(blank_layout)
    # 整页深红背景
    add_rect(slide, Emu(0), Emu(0), SLIDE_W, SLIDE_H, COLOR_PRIMARY)
    # 顶部金色装饰条
    add_rect(slide, Emu(0), Inches(0.8), SLIDE_W, Inches(0.04), COLOR_ACCENT)
    # 底部金色装饰条
    add_rect(slide, Emu(0), SLIDE_H - Inches(0.8),
             SLIDE_W, Inches(0.04), COLOR_ACCENT)
    # 右上小标识
    add_text_box(slide, SLIDE_W - Inches(4), Inches(1.0),
                 Inches(3.6), Inches(0.4), "民航系统党课  ·  专题汇报",
                 size=14, color=COLOR_ACCENT, align=PP_ALIGN.RIGHT,
                 font=FONT_BODY)

    # 主标题
    add_text_box(slide, Inches(1.0), Inches(2.4), Inches(11.3),
                 Inches(1.2),
                 "坚守正确政绩观",
                 size=60, bold=True, color=COLOR_WHITE,
                 align=PP_ALIGN.CENTER, font=FONT_TITLE)
    add_text_box(slide, Inches(1.0), Inches(3.6), Inches(11.3),
                 Inches(1.2),
                 "筑牢民航初心使命",
                 size=60, bold=True, color=COLOR_ACCENT,
                 align=PP_ALIGN.CENTER, font=FONT_TITLE)

    # 副标题
    add_text_box(slide, Inches(1.0), Inches(5.0), Inches(11.3),
                 Inches(0.6),
                 "— 民航运输领域典型案例党课 —",
                 size=20, color=COLOR_WHITE,
                 align=PP_ALIGN.CENTER, font=FONT_BODY)

    # 落款
    add_text_box(slide, Inches(1.0), SLIDE_H - Inches(1.5),
                 Inches(11.3), Inches(0.4),
                 "汇报单位：民航党支部",
                 size=14, color=COLOR_WHITE,
                 align=PP_ALIGN.CENTER, font=FONT_BODY)
    add_text_box(slide, Inches(1.0), SLIDE_H - Inches(1.1),
                 Inches(11.3), Inches(0.4),
                 "汇 报 人：×××",
                 size=14, color=COLOR_WHITE,
                 align=PP_ALIGN.CENTER, font=FONT_BODY)


# ---------------- Slide 2: 目录 ----------------
def build_toc():
    slide = prs.slides.add_slide(blank_layout)
    add_slide_header(slide, "目  录  CONTENTS",
                     "本次党课主要从以下四个方面展开")

    items = [
        ("01", "深刻把握民航运输领域正确政绩观核心要义",
         "明确内涵 · 厘清边界"),
        ("02", "以厦航为标杆，践行运输领域正确政绩观",
         "正面示范 · 经验做法"),
        ("03", "负面典型：民航运输领域政绩观错位事例",
         "五大案例 · 警钟长鸣"),
        ("04", "正反对照反思，校准民航干部正确政绩观",
         "对标对表 · 担当作为"),
    ]
    top0 = Inches(1.9)
    row_h = Inches(1.15)
    for i, (num, title, sub) in enumerate(items):
        y = top0 + row_h * i
        # 序号方块
        add_rect(slide, Inches(0.8), y, Inches(1.0), Inches(0.95),
                 COLOR_PRIMARY)
        add_text_box(slide, Inches(0.8), y, Inches(1.0), Inches(0.95),
                     num, size=34, bold=True, color=COLOR_ACCENT,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                     font=FONT_TITLE)
        # 标题
        add_text_box(slide, Inches(2.1), y + Inches(0.08),
                     Inches(10.5), Inches(0.5), title,
                     size=20, bold=True, color=COLOR_DARK,
                     anchor=MSO_ANCHOR.MIDDLE, font=FONT_TITLE)
        # 副标题
        add_text_box(slide, Inches(2.1), y + Inches(0.55),
                     Inches(10.5), Inches(0.4), sub,
                     size=13, color=COLOR_TEXT_LIGHT,
                     anchor=MSO_ANCHOR.MIDDLE, font=FONT_BODY)
        # 分隔线
        if i < len(items) - 1:
            add_rect(slide, Inches(2.1), y + row_h - Inches(0.05),
                     Inches(10.5), Emu(8000), COLOR_BG_LIGHT)

    add_page_footer(slide, 2, TOTAL_PAGES, "目  录")


# ---------------- Slide 3: 导语 ----------------
def build_preface():
    slide = prs.slides.add_slide(blank_layout)
    add_slide_header(slide, "导  语  PREFACE",
                     "政绩观是党员干部世界观、人生观、价值观的集中体现")

    # 主要文本
    text_para = (
        "同志们，政绩观是党员干部世界观、人生观、价值观的集中体现，"
        "是干事创业、履职尽责的根本标尺。民航作为国家战略性先导产业，"
        "一头连着国家安全、一头连着民生福祉，运输领域的航线布局、"
        "市场经营、服务保障、行风作风、国企管理等每一项工作，"
        "都直接检验我们的政绩观。"
    )
    add_text_box(slide, Inches(0.8), Inches(1.9), Inches(11.7),
                 Inches(2.0), text_para,
                 size=18, color=COLOR_TEXT,
                 line_spacing=1.6, font=FONT_BODY)

    # 三大要点卡片
    cards = [
        ("聚焦", "运输主业政绩观",
         "经营、市场、航线、服务、管理"),
        ("剖析", "五类反面问题",
         "错位 · 形式 · 功利 · 漂浮 · 失范"),
        ("落脚", "正反对照以案明纪",
         "以案促改 · 担当作为"),
    ]
    card_w = Inches(3.8)
    card_h = Inches(2.2)
    gap = Inches(0.25)
    total_w = card_w * 3 + gap * 2
    start_x = (SLIDE_W - total_w) / 2
    y = Inches(4.4)
    for i, (tag, t1, t2) in enumerate(cards):
        x = start_x + (card_w + gap) * i
        add_round_rect(slide, x, y, card_w, card_h, COLOR_BG_CARD)
        # 顶部色条
        add_rect(slide, x, y, card_w, Inches(0.45), COLOR_PRIMARY)
        add_text_box(slide, x, y, card_w, Inches(0.45),
                     tag, size=18, bold=True, color=COLOR_ACCENT,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                     font=FONT_TITLE)
        add_text_box(slide, x, y + Inches(0.65), card_w, Inches(0.6),
                     t1, size=19, bold=True, color=COLOR_DARK,
                     align=PP_ALIGN.CENTER, font=FONT_TITLE)
        add_text_box(slide, x, y + Inches(1.35), card_w, Inches(0.7),
                     t2, size=13, color=COLOR_TEXT,
                     align=PP_ALIGN.CENTER, font=FONT_BODY,
                     line_spacing=1.4)

    add_page_footer(slide, 3, TOTAL_PAGES, "导  语")


# ---------------- Slide 4: 第一章封面 ----------------
def build_section_cover(num, title, subtitle, page_num, section_name):
    slide = prs.slides.add_slide(blank_layout)
    # 左侧深红块
    add_rect(slide, Emu(0), Emu(0), Inches(5.0), SLIDE_H, COLOR_PRIMARY)
    # 右侧白色背景已是默认
    # 大序号
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(4.0), Inches(3.0),
                 num, size=200, bold=True, color=COLOR_ACCENT,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 font=FONT_TITLE)
    # 装饰横线
    add_rect(slide, Inches(5.5), Inches(2.0), Inches(0.7), Emu(40000),
             COLOR_PRIMARY)
    # PART 标识
    add_text_box(slide, Inches(5.5), Inches(2.2), Inches(3), Inches(0.4),
                 f"PART  {num}", size=16, bold=True, color=COLOR_PRIMARY,
                 font=FONT_TITLE)
    # 标题
    add_text_box(slide, Inches(5.5), Inches(2.9), Inches(7.5), Inches(1.5),
                 title, size=30, bold=True, color=COLOR_DARK,
                 line_spacing=1.3, font=FONT_TITLE)
    # 副标题
    add_text_box(slide, Inches(5.5), Inches(5.0), Inches(7.5), Inches(1.0),
                 subtitle, size=14, color=COLOR_TEXT_LIGHT,
                 line_spacing=1.5, font=FONT_BODY)

    add_page_footer(slide, page_num, TOTAL_PAGES, section_name)


# ---------------- Slide 5: 第一章内容 ----------------
def build_part1_content():
    slide = prs.slides.add_slide(blank_layout)
    add_slide_header(slide, "一、正确政绩观核心要义",
                     "对民航运输系统党员干部而言，正确政绩观就是——")

    # 三句话核心
    add_round_rect(slide, Inches(0.8), Inches(1.8),
                   Inches(11.7), Inches(1.4), COLOR_BG_CARD)
    add_text_box(slide, Inches(1.0), Inches(1.95), Inches(11.3),
                 Inches(0.6),
                 "坚持人民航空为人民",
                 size=24, bold=True, color=COLOR_PRIMARY,
                 align=PP_ALIGN.CENTER, font=FONT_TITLE)
    add_text_box(slide, Inches(1.0), Inches(2.55), Inches(11.3),
                 Inches(0.6),
                 "不唯数据唯实效 · 不重显绩重潜绩 · 不谋私利谋公利",
                 size=18, color=COLOR_DARK,
                 align=PP_ALIGN.CENTER, font=FONT_BODY)

    # 左右两栏 倡导 vs 反对
    # 左栏：倡导
    left_x = Inches(0.8)
    col_w = Inches(5.85)
    col_y = Inches(3.5)
    col_h = Inches(3.2)

    add_rect(slide, left_x, col_y, col_w, Inches(0.5), COLOR_GREEN)
    add_text_box(slide, left_x, col_y, col_w, Inches(0.5),
                 "✓  坚持倡导",
                 size=18, bold=True, color=COLOR_WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 font=FONT_TITLE)
    add_rect(slide, left_x, col_y + Inches(0.5), col_w, Inches(2.7),
             COLOR_BG_CARD)
    good_items = [
        "实事求是、务实为民",
        "合规经营、久久为功",
        "深耕主业、稳健发展",
        "以人民为中心、以旅客为本",
    ]
    paras = []
    for it in good_items:
        paras.append({"text": "●  " + it, "size": 16,
                      "color": COLOR_TEXT, "bold": False,
                      "space_after": 8})
    add_multi_paragraphs(slide, left_x + Inches(0.4),
                         col_y + Inches(0.7), col_w - Inches(0.6),
                         Inches(2.4), paras, line_spacing=1.5)

    # 右栏：反对
    right_x = Inches(6.85)
    add_rect(slide, right_x, col_y, col_w, Inches(0.5), COLOR_RED_NEG)
    add_text_box(slide, right_x, col_y, col_w, Inches(0.5),
                 "✗  坚决反对",
                 size=18, bold=True, color=COLOR_WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 font=FONT_TITLE)
    add_rect(slide, right_x, col_y + Inches(0.5), col_w, Inches(2.7),
             COLOR_BG_CARD)
    bad_items = [
        "盲目铺摊子、虚假凑指标",
        "重宣传轻实效、搞行业特权",
        "经营粗放任性、违规失范",
        "不顾基层和群众实际",
    ]
    paras = []
    for it in bad_items:
        paras.append({"text": "●  " + it, "size": 16,
                      "color": COLOR_TEXT, "bold": False,
                      "space_after": 8})
    add_multi_paragraphs(slide, right_x + Inches(0.4),
                         col_y + Inches(0.7), col_w - Inches(0.6),
                         Inches(2.4), paras, line_spacing=1.5)

    add_page_footer(slide, 5, TOTAL_PAGES,
                    "第一部分  正确政绩观核心要义")


# ---------------- Slide 7: 厦航党委班子四点示范 ----------------
def build_xiamen_committee():
    slide = prs.slides.add_slide(blank_layout)
    add_slide_header(slide, "（一）厦航党委班子：锚定运输主业，一张蓝图绘到底",
                     "依靠改革 · 舍得投入 · 服务规范 · 以人为本")

    items = [
        ("01", "稳健深耕运输市场",
         "立足福建区位优势，深耕干线、做精支线、做强两岸航线网络，"
         "做一条、稳一条、优一条，不贪大求洋、不搞短期冲量。"),
        ("02", "把服务做成口碑政绩",
         "以旅客出行体验为出发点，优化中转衔接、完善通程航班、"
         "升级地面服务，靠真实口碑获评中国质量奖、世界级航空公司。"),
        ("03", "关键时刻扛起政治责任",
         "疫情期间主动承担援鄂援沪包机、医疗物资转运、旅客疏运任务；"
         "重大节假日和极端天气优先保障旅客出行，彰显国企担当。"),
        ("04", "坚守合规涵养清风正气",
         "严格规范航线、运价、代理、采购等关键领域；班子不搞特权、"
         "不插手项目、不利用资源谋私利，制度管人、流程管事。"),
    ]
    # 2x2 网格
    card_w = Inches(6.0)
    card_h = Inches(2.4)
    gap_x = Inches(0.3)
    gap_y = Inches(0.2)
    start_x = (SLIDE_W - card_w * 2 - gap_x) / 2
    start_y = Inches(1.85)
    for i, (num, title, desc) in enumerate(items):
        col = i % 2
        row = i // 2
        x = start_x + (card_w + gap_x) * col
        y = start_y + (card_h + gap_y) * row
        add_round_rect(slide, x, y, card_w, card_h, COLOR_BG_CARD)
        # 序号方块
        add_rect(slide, x + Inches(0.2), y + Inches(0.25),
                 Inches(0.85), Inches(0.85), COLOR_PRIMARY)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.25),
                     Inches(0.85), Inches(0.85), num,
                     size=22, bold=True, color=COLOR_ACCENT,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                     font=FONT_TITLE)
        # 标题
        add_text_box(slide, x + Inches(1.2), y + Inches(0.3),
                     card_w - Inches(1.4), Inches(0.55),
                     title, size=17, bold=True, color=COLOR_DARK,
                     font=FONT_TITLE, anchor=MSO_ANCHOR.MIDDLE)
        # 描述
        add_text_box(slide, x + Inches(0.25), y + Inches(1.15),
                     card_w - Inches(0.5), card_h - Inches(1.25),
                     desc, size=13, color=COLOR_TEXT,
                     line_spacing=1.5, font=FONT_BODY)

    add_page_footer(slide, 7, TOTAL_PAGES,
                    "第二部分  以厦航为标杆")


# ---------------- Slide 8: 厦航历任领导 ----------------
def build_xiamen_leaders():
    slide = prs.slides.add_slide(blank_layout)
    add_slide_header(slide, "（二）厦航历任领导示范引领",
                     "实干兴航 · 稳健经营 · 服务大局")

    leaders = [
        ("吴荣南", "创业奠基",
         "白手起家，不靠政策等靠要",
         "脚踏实地拓航线、稳客源，与员工同甘共苦，不搞特殊待遇，"
         "奠定厦航务实低调、实干兴航的底色。"),
        ("车尚轮", "品质深耕",
         "质量优先，稳健经营不虚胖",
         "深耕运输品质、打磨服务链条，以长期主义做运输、做品牌，"
         "不靠短期数据注水刷政绩。"),
        ("赵  东", "融入大局",
         "服务地方，融入区域大局",
         "优化省内干支航线网络，助力乡村振兴和文旅发展，做强两岸"
         "航空运输通道，把企业发展融入国家和区域大局。"),
    ]
    card_w = Inches(4.0)
    card_h = Inches(4.4)
    gap = Inches(0.25)
    total_w = card_w * 3 + gap * 2
    start_x = (SLIDE_W - total_w) / 2
    y = Inches(1.9)
    for i, (name, tag, headline, desc) in enumerate(leaders):
        x = start_x + (card_w + gap) * i
        # 卡片底
        add_round_rect(slide, x, y, card_w, card_h, COLOR_BG_CARD)
        # 顶部红条
        add_rect(slide, x, y, card_w, Inches(1.4), COLOR_PRIMARY)
        # 姓名
        add_text_box(slide, x, y + Inches(0.2), card_w, Inches(0.65),
                     name, size=30, bold=True, color=COLOR_WHITE,
                     align=PP_ALIGN.CENTER, font=FONT_TITLE)
        # 标签
        add_text_box(slide, x, y + Inches(0.9), card_w, Inches(0.4),
                     tag, size=14, color=COLOR_ACCENT,
                     align=PP_ALIGN.CENTER, font=FONT_BODY)
        # 一句话
        add_text_box(slide, x + Inches(0.25), y + Inches(1.65),
                     card_w - Inches(0.5), Inches(0.8),
                     headline, size=16, bold=True, color=COLOR_DARK,
                     align=PP_ALIGN.CENTER, font=FONT_TITLE,
                     line_spacing=1.4)
        # 分隔线
        add_rect(slide, x + Inches(0.7), y + Inches(2.5),
                 card_w - Inches(1.4), Emu(8000), COLOR_ACCENT)
        # 描述
        add_text_box(slide, x + Inches(0.3), y + Inches(2.7),
                     card_w - Inches(0.6), Inches(1.6),
                     desc, size=13, color=COLOR_TEXT,
                     align=PP_ALIGN.LEFT, line_spacing=1.6,
                     font=FONT_BODY)

    add_page_footer(slide, 8, TOTAL_PAGES,
                    "第二部分  以厦航为标杆")


# ---------------- 反面事例通用模板 ----------------
def build_negative_case(idx, title, tag, problem, harm, page_num):
    slide = prs.slides.add_slide(blank_layout)
    add_slide_header(slide, f"反面事例 {idx}：{title}",
                     "重显绩轻潜绩 · 重形式轻实效 · 政绩观错位")

    # 警示标签
    add_round_rect(slide, Inches(0.8), Inches(1.8),
                   Inches(2.5), Inches(0.55), COLOR_RED_NEG)
    add_text_box(slide, Inches(0.8), Inches(1.8),
                 Inches(2.5), Inches(0.55),
                 f"⚠  {tag}", size=16, bold=True,
                 color=COLOR_WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE, font=FONT_TITLE)

    # 问题表现
    y1 = Inches(2.65)
    add_rect(slide, Inches(0.8), y1, Inches(0.15), Inches(0.5),
             COLOR_PRIMARY)
    add_text_box(slide, Inches(1.05), y1, Inches(11), Inches(0.5),
                 "问题表现", size=18, bold=True, color=COLOR_PRIMARY,
                 anchor=MSO_ANCHOR.MIDDLE, font=FONT_TITLE)
    add_round_rect(slide, Inches(0.8), y1 + Inches(0.55),
                   Inches(11.7), Inches(2.0), COLOR_BG_CARD)
    add_text_box(slide, Inches(1.1), y1 + Inches(0.7),
                 Inches(11.1), Inches(1.8),
                 problem, size=14, color=COLOR_TEXT,
                 line_spacing=1.6, font=FONT_BODY)

    # 危害剖析
    y2 = Inches(5.4)
    add_rect(slide, Inches(0.8), y2, Inches(0.15), Inches(0.5),
             COLOR_RED_NEG)
    add_text_box(slide, Inches(1.05), y2, Inches(11), Inches(0.5),
                 "本质剖析", size=18, bold=True, color=COLOR_RED_NEG,
                 anchor=MSO_ANCHOR.MIDDLE, font=FONT_TITLE)
    add_text_box(slide, Inches(0.85), y2 + Inches(0.55),
                 Inches(11.6), Inches(1.0),
                 harm, size=14, bold=True, color=COLOR_DARK,
                 line_spacing=1.5, font=FONT_BODY)

    add_page_footer(slide, page_num, TOTAL_PAGES,
                    "第三部分  反面典型事例")


# ---------------- 反思页 ----------------
def build_reflection():
    slide = prs.slides.add_slide(blank_layout)
    add_slide_header(slide, "四、正反对照反思  校准正确政绩观",
                     "对照厦航实干稳健、为民务实、合规清廉的正面标杆")

    items = [
        ("01", "坚决不搞形象航线、不做数据政绩",
         "一切工作以旅客真实需求、市场实际规律、企业长远发展为标准"),
        ("02", "坚决不利用航线时刻、运输资源搞特权谋私利",
         "守住公权为公、资源为民的底线"),
        ("03", "坚决摒弃形式主义，少做表面文章、多解民生实事",
         "把运输服务的功夫下在群众心坎上"),
        ("04", "坚持稳健经营、久久为功",
         "一任接着一任干，创造经得起实践、群众和历史检验的实绩"),
    ]
    y0 = Inches(1.85)
    row_h = Inches(1.18)
    for i, (num, headline, desc) in enumerate(items):
        y = y0 + row_h * i
        # 序号圆形
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                        Inches(0.8), y + Inches(0.1),
                                        Inches(0.95), Inches(0.95))
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLOR_PRIMARY
        circle.line.fill.background()
        circle.shadow.inherit = False
        circle.text_frame.text = ""
        add_text_box(slide, Inches(0.8), y + Inches(0.1),
                     Inches(0.95), Inches(0.95),
                     num, size=22, bold=True, color=COLOR_ACCENT,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                     font=FONT_TITLE)
        # 标题
        add_text_box(slide, Inches(2.0), y + Inches(0.1),
                     Inches(10.5), Inches(0.55),
                     headline, size=17, bold=True, color=COLOR_DARK,
                     anchor=MSO_ANCHOR.MIDDLE, font=FONT_TITLE)
        # 描述
        add_text_box(slide, Inches(2.0), y + Inches(0.62),
                     Inches(10.5), Inches(0.5),
                     desc, size=13, color=COLOR_TEXT_LIGHT,
                     font=FONT_BODY)
        # 分隔线
        if i < len(items) - 1:
            add_rect(slide, Inches(2.0), y + row_h - Inches(0.05),
                     Inches(10.5), Emu(8000), COLOR_BG_LIGHT)

    add_page_footer(slide, 15, TOTAL_PAGES,
                    "第四部分  反思与校准")


# ---------------- 三点要求页 ----------------
def build_three_requirements():
    slide = prs.slides.add_slide(blank_layout)
    add_slide_header(slide, "党员干部三点要求  立足本职担当作为",
                     "作为民航系统党员干部")

    items = [
        ("一", "对标厦航",
         "学习其务实经营、服务为民、稳健长远的发展理念，"
         "把先进经验转化为本职工作思路与方法。"),
        ("二", "以案为戒",
         "从运输领域反面典型中受警醒、明底线、知敬畏，"
         "时刻校准政绩观、权力观、利益观。"),
        ("三", "立足本职",
         "在航线保障、客运服务、市场经营、行风建设中坚守初心，"
         "为民航运输高质量发展履职尽责、担当作为。"),
    ]
    card_w = Inches(4.0)
    card_h = Inches(4.6)
    gap = Inches(0.25)
    total_w = card_w * 3 + gap * 2
    start_x = (SLIDE_W - total_w) / 2
    y = Inches(1.85)
    for i, (num, title, desc) in enumerate(items):
        x = start_x + (card_w + gap) * i
        # 主卡片
        add_round_rect(slide, x, y, card_w, card_h, COLOR_BG_CARD)
        # 大数字
        add_text_box(slide, x, y + Inches(0.3),
                     card_w, Inches(1.4),
                     num, size=64, bold=True, color=COLOR_PRIMARY,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                     font=FONT_TITLE)
        # 分隔线
        add_rect(slide, x + Inches(1.3), y + Inches(1.9),
                 card_w - Inches(2.6), Emu(20000), COLOR_ACCENT)
        # 标题
        add_text_box(slide, x, y + Inches(2.05),
                     card_w, Inches(0.7),
                     title, size=24, bold=True, color=COLOR_DARK,
                     align=PP_ALIGN.CENTER, font=FONT_TITLE)
        # 描述
        add_text_box(slide, x + Inches(0.4),
                     y + Inches(2.85),
                     card_w - Inches(0.8), Inches(1.6),
                     desc, size=14, color=COLOR_TEXT,
                     align=PP_ALIGN.LEFT, line_spacing=1.6,
                     font=FONT_BODY)

    add_page_footer(slide, 16, TOTAL_PAGES,
                    "第四部分  反思与校准")


# ---------------- 结语页 ----------------
def build_closing():
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, Emu(0), Emu(0), SLIDE_W, SLIDE_H, COLOR_PRIMARY)
    add_rect(slide, Emu(0), Inches(0.8), SLIDE_W, Inches(0.04), COLOR_ACCENT)
    add_rect(slide, Emu(0), SLIDE_H - Inches(0.8), SLIDE_W,
             Inches(0.04), COLOR_ACCENT)

    add_text_box(slide, Inches(1.0), Inches(1.8), Inches(11.3),
                 Inches(0.6), "结  语", size=24, color=COLOR_ACCENT,
                 align=PP_ALIGN.CENTER, font=FONT_TITLE)

    add_text_box(slide, Inches(1.0), Inches(2.6), Inches(11.3),
                 Inches(1.2), "树立和践行正确政绩观",
                 size=48, bold=True, color=COLOR_WHITE,
                 align=PP_ALIGN.CENTER, font=FONT_TITLE)
    add_text_box(slide, Inches(1.0), Inches(3.7), Inches(11.3),
                 Inches(1.2), "为民航高质量发展担当作为",
                 size=48, bold=True, color=COLOR_ACCENT,
                 align=PP_ALIGN.CENTER, font=FONT_TITLE)

    # 装饰横线
    add_rect(slide, Inches(5.5), Inches(5.2), Inches(2.3),
             Emu(20000), COLOR_ACCENT)

    add_text_box(slide, Inches(1.0), Inches(5.5), Inches(11.3),
                 Inches(0.6),
                 "履职尽责  ·  实干兴航  ·  久久为功",
                 size=22, color=COLOR_WHITE,
                 align=PP_ALIGN.CENTER, font=FONT_BODY)

    add_text_box(slide, Inches(1.0), SLIDE_H - Inches(1.2),
                 Inches(11.3), Inches(0.6),
                 "感谢聆听  恳请批评指正",
                 size=20, color=COLOR_WHITE,
                 align=PP_ALIGN.CENTER, font=FONT_BODY)


# ============ 按顺序构建 ============
build_cover()                # 1
build_toc()                  # 2
build_preface()              # 3

build_section_cover("01",
                    "深刻把握民航运输领域\n正确政绩观核心要义",
                    "Understanding the Core Meaning\n"
                    "明确正确政绩观的内涵与边界，"
                    "厘清倡导与反对",
                    4,
                    "第一部分  正确政绩观核心要义")  # 4
build_part1_content()  # 5

build_section_cover("02",
                    "以厦航为标杆\n践行运输领域正确政绩观",
                    "Xiamen Airlines as a Benchmark\n"
                    "锚定运输主业，一张蓝图绘到底",
                    6,
                    "第二部分  以厦航为标杆")  # 6
build_xiamen_committee()  # 7
build_xiamen_leaders()    # 8

build_section_cover("03",
                    "负面典型\n民航运输领域政绩观错位事例",
                    "Negative Cases for Warning\n"
                    "五类反面问题，警钟长鸣",
                    9,
                    "第三部分  反面典型事例")  # 9

build_negative_case(
    1, "航线开发搞\"形象工程\"", "形式主义",
    "部分地区民航管理部门和支线机场、航企干部，为完成航线开通考核指标、"
    "应付上级检查、打造宣传亮点，盲目申报新开航线，不顾客流基础、"
    "市场支撑、收益测算，只为凑航线数量、做开通仪式、上新闻报道；"
    "航线开通后不做市场培育、不做客源引流、不优化时刻班次，"
    "短期内客座率低迷、持续大额亏损，随后被迫停飞。",
    "只求\"开通即政绩\"，不求\"长久惠民生\"，"
    "浪费运力资源、财政补贴资源，典型的重显绩轻潜绩、重形式轻实效。",
    10)  # 10

build_negative_case(
    2, "运输数据造假玩\"数字游戏\"", "弄虚作假",
    "个别航企、地方民航单位，为冲旅客吞吐量、航班量、中转量、"
    "市场份额等考核数据：人为拆分航班班次、虚增中转旅客统计、"
    "内部调账美化经营数据；搞虚假合作、虚构客源协议，报表好看、"
    "实际经营惨淡；把精力放在\"修饰数据、应付考核\"上，"
    "不深耕市场、不优化服务、不挖掘真实客源。",
    "误导行业规划和资源投放，典型的弄虚作假、唯指标论的错误政绩观，"
    "严重损害民航行业治理基础和公信力。",
    11)  # 11

build_negative_case(
    3, "利用运输资源搞特权谋私利", "权力寻租",
    "民航运输领域航线审批、时刻分配、客运代理、运价管控、包机业务、"
    "航班合作等权力集中领域，少数干部把公共航空运输资源当作私人筹码，"
    "插手航线申请、时刻调整，为关联企业、熟人关系倾斜资源；"
    "在客运销售、代理招商、合作外包中优亲厚友、吃拿卡要，"
    "破坏运输市场公平秩序。",
    "信奉\"有权就有特权\"，把行业公共资源异化为个人谋利工具，"
    "破坏市场公平秩序，背离为民服务初心。",
    12)  # 12

build_negative_case(
    4, "服务作风漂浮搞表面文章", "作风漂浮",
    "部分民航运输一线管理干部，抓服务只做表面文章：热衷搞硬件装修、"
    "打造网红候机区域、做宣传展板和宣传片，却不解决旅客反映强烈的"
    "中转不便、票务退改繁琐、航班信息不透明、地面衔接不畅等实际问题；"
    "检查督导只看台账、看现场布置、看拍照留痕，"
    "不深入旅客、不倾听一线员工心声。",
    "工作停留在\"做了、说了、宣传了\"，没有做到\"做实、做好、群众满意\"，"
    "是典型重表面、轻实效的形式主义政绩观。",
    13)  # 13

build_negative_case(
    5, "国企经营粗放盲目扩张", "短期功利",
    "个别民航运输类国企班子，政绩观短期化、功利化：不顾自身经营实力，"
    "盲目引进大量飞机、扩张运力、铺设过多网点；粗放管理、成本管控缺失，"
    "盲目拼价格、打恶性市场价格战，扰乱行业运输秩序；"
    "只追求任期内规模做大、场面做亮，留下高额负债、低效资产、"
    "经营隐患给下一任。",
    "缺乏功成不必在我的长远担当，典型的寅吃卯粮、只顾眼前的错误政绩观，"
    "贻害企业可持续发展和行业健康生态。",
    14)  # 14

build_reflection()           # 15
build_three_requirements()   # 16
build_closing()              # 17

# 实际总页数
real_total = len(prs.slides)
print(f"生成幻灯片总数: {real_total}")

output_path = "/home/user/cello-practise/presentation/坚守正确政绩观-民航党课汇报.pptx"
prs.save(output_path)
print(f"已保存到: {output_path}")
